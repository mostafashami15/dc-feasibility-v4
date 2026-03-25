"""
Generate DC Feasibility Tool v4 presentation slide.
Matches the Metlen team deck style (colors, fonts, layout density).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── Colors (from Metlen theme) ──────────────────────────────────────
NAVY       = RGBColor(0x0A, 0x22, 0x40)
PURPLE     = RGBColor(0x79, 0x5A, 0xFD)
DARK_PURPLE= RGBColor(0x4E, 0x25, 0x89)
CYAN       = RGBColor(0x00, 0xF1, 0xF2)
GREEN      = RGBColor(0x5F, 0xE8, 0x38)
LAVENDER   = RGBColor(0xBC, 0xAC, 0xFE)
DARK_TEAL  = RGBColor(0x1B, 0x30, 0x3A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)
BODY_GRAY  = RGBColor(0x4A, 0x4A, 0x4A)

# ── Dimensions ───────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def emu(inches):
    return Inches(inches)

def add_gradient_fill(shape, color1, color2):
    """Add a linear gradient fill to a shape."""
    spPr = shape._element.spPr
    # Remove existing fill
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'noFill'):
            spPr.remove(child)

    gradFill = spPr.makeelement(qn('a:gradFill'), {})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})

    gs1 = gsLst.makeelement(qn('a:gs'), {'pos': '0'})
    srgb1 = gs1.makeelement(qn('a:srgbClr'), {'val': str(color1)})
    gs1.append(srgb1)
    gsLst.append(gs1)

    gs2 = gsLst.makeelement(qn('a:gs'), {'pos': '100000'})
    srgb2 = gs2.makeelement(qn('a:srgbClr'), {'val': str(color2)})
    gs2.append(srgb2)
    gsLst.append(gs2)

    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn('a:lin'), {'ang': '0', 'scaled': '1'})
    gradFill.append(lin)
    spPr.append(gradFill)

def set_shape_border(shape, color, width_pt=1):
    """Set shape outline."""
    ln = shape._element.spPr.makeelement(qn('a:ln'), {'w': str(int(width_pt * 12700))})
    solidFill = ln.makeelement(qn('a:solidFill'), {})
    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': str(color)})
    solidFill.append(srgb)
    ln.append(solidFill)
    # Remove existing ln
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}ln') or child.tag == 'ln':
            spPr.remove(child)
    spPr.append(ln)

def no_border(shape):
    """Remove shape outline."""
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}ln') or child.tag == 'ln':
            spPr.remove(child)
    ln = spPr.makeelement(qn('a:ln'), {})
    noFill = ln.makeelement(qn('a:noFill'), {})
    ln.append(noFill)
    spPr.append(ln)

def add_textbox(slide, left, top, width, height, text, font_name='Open Sans',
                font_size=Pt(10), bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))

    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_name='Open Sans',
                          font_size=Pt(9), color=BLACK, line_spacing=Pt(14),
                          bold_first=False, alignment=PP_ALIGN.LEFT, bullet_color=None):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(1)
        p.space_after = Pt(1)

        # Check if line has a bullet prefix
        if line.startswith(('\u2022 ', '- ')):
            if bullet_color:
                bullet_run = p.add_run()
                bullet_run.text = line[0] + ' '
                bullet_run.font.name = font_name
                bullet_run.font.size = font_size
                bullet_run.font.color.rgb = bullet_color
                bullet_run.font.bold = True

                text_run = p.add_run()
                text_run.text = line[2:]
                text_run.font.name = font_name
                text_run.font.size = font_size
                text_run.font.color.rgb = color
                text_run.font.bold = False
                continue

        # Check for bold:regular split with |
        if '|' in line and not line.startswith('|'):
            parts = line.split('|', 1)
            r1 = p.add_run()
            r1.text = parts[0]
            r1.font.name = font_name
            r1.font.size = font_size
            r1.font.color.rgb = color
            r1.font.bold = True

            r2 = p.add_run()
            r2.text = parts[1]
            r2.font.name = font_name
            r2.font.size = Pt(max(font_size.pt - 1, 7))
            r2.font.color.rgb = BODY_GRAY
            r2.font.bold = False
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = (bold_first and i == 0)

    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=1):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        set_shape_border(shape, border_color, border_width)
    else:
        no_border(shape)
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape

def add_pentagon(slide, left, top, width, height, fill_color1, fill_color2, text='',
                 font_size=Pt(9), font_color=WHITE, font_name='Montserrat SemiBold'):
    """Add a pentagon (chevron arrow) with gradient and text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    add_gradient_fill(shape, fill_color1, fill_color2)
    no_border(shape)

    # Reduce the notch depth
    if shape.adjustments:
        shape.adjustments[0] = 0.25

    if text:
        shape.text_frame.word_wrap = True
        shape.text_frame.auto_size = None
        shape.text_frame._txBody.bodyPr.set('anchor', 'ctr')
        # Set left/right margins to accommodate chevron shape
        shape.text_frame._txBody.bodyPr.set('lIns', str(int(Inches(0.25))))
        shape.text_frame._txBody.bodyPr.set('rIns', str(int(Inches(0.05))))

        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = True

    return shape

def add_arrow_connector(slide, x1, y1, x2, y2, color=PURPLE, width_pt=1.5):
    """Add a line connector."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    ln = connector._element.spPr.makeelement(qn('a:ln'), {'w': str(int(width_pt * 12700))})
    solidFill = ln.makeelement(qn('a:solidFill'), {})
    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': str(color)})
    solidFill.append(srgb)
    ln.append(solidFill)

    # Add arrow head
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)

    spPr = connector._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}ln'):
            spPr.remove(child)
    spPr.append(ln)
    return connector


def build_slide():
    """Build the DC Feasibility Tool v4 slide."""
    # Load team PPTX to get theme
    team_path = os.path.join(os.path.dirname(__file__), 'Metlen - Data Center BU Mapping.pptx')
    prs = Presentation(team_path)

    # Use existing slide dimensions
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Add blank slide
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower() or 'Blank' in layout.name:
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank_layout)

    # Remove any placeholder shapes from the slide
    for shape in list(slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    # ════════════════════════════════════════════════════════════════
    # HEADER AREA
    # ════════════════════════════════════════════════════════════════

    # Top colored bar (navy)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    no_border(top_bar)

    # Pentagon label "Feasibility Analysis"
    pent = add_pentagon(slide, Inches(0.4), Inches(0.25), Inches(2.8), Inches(0.42),
                        PURPLE, DARK_PURPLE, 'Feasibility Analysis Tool',
                        font_size=Pt(11), font_color=WHITE)

    # Title
    add_textbox(slide, Inches(3.4), Inches(0.15), Inches(7), Inches(0.45),
                'DC Feasibility Tool v4', font_name='Manrope Light',
                font_size=Pt(24), bold=False, color=DARK_PURPLE)

    # Subtitle
    add_textbox(slide, Inches(3.4), Inches(0.52), Inches(8), Inches(0.3),
                'Site-specific 8,760-hour simulation engine for data center site evaluation',
                font_name='Open Sans', font_size=Pt(11), bold=False, color=BODY_GRAY)

    # Thin separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.015))
    sep.fill.solid()
    sep.fill.fore_color.rgb = LAVENDER
    no_border(sep)

    # ════════════════════════════════════════════════════════════════
    # WORKFLOW ROW (5 pentagons)
    # ════════════════════════════════════════════════════════════════

    workflow_y = Inches(1.1)
    workflow_h = Inches(0.55)
    pent_w = Inches(2.2)
    gap = Inches(0.18)
    start_x = Inches(0.7)

    steps = [
        ('1  Site & Grid', DARK_PURPLE, PURPLE),
        ('2  Climate Analysis', PURPLE, LAVENDER),
        ('3  Scenario Runner', DARK_PURPLE, PURPLE),
        ('4  Results Dashboard', PURPLE, LAVENDER),
        ('5  Green Energy', DARK_PURPLE, PURPLE),
    ]

    for i, (label, c1, c2) in enumerate(steps):
        x = start_x + i * (pent_w + gap)
        add_pentagon(slide, x, workflow_y, pent_w, workflow_h, c1, c2, label,
                     font_size=Pt(10), font_color=WHITE)

    # "Workflow" label
    add_textbox(slide, Inches(0.4), Inches(1.12), Inches(0.3), Inches(0.5),
                '', font_size=Pt(8))

    # ════════════════════════════════════════════════════════════════
    # THREE-COLUMN CONTENT AREA
    # ════════════════════════════════════════════════════════════════

    col_y = Inches(1.85)
    col_h = Inches(3.65)
    margin = Inches(0.4)
    col_gap = Inches(0.25)

    # Column widths
    col1_w = Inches(3.9)
    col2_w = Inches(4.4)
    col3_w = Inches(3.5)

    col1_x = margin
    col2_x = col1_x + col1_w + col_gap
    col3_x = col2_x + col2_w + col_gap

    # ── Column 1: DATA SOURCES ───────────────────────────────────
    box1 = add_rounded_rect(slide, col1_x, col_y, col1_w, col_h,
                            fill_color=LIGHT_GRAY, border_color=LAVENDER, border_width=1)

    # Column header
    hdr1_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_x, col_y, col1_w, Inches(0.35))
    hdr1_bg.fill.solid()
    hdr1_bg.fill.fore_color.rgb = NAVY
    no_border(hdr1_bg)
    add_textbox(slide, col1_x + Inches(0.15), col_y + Inches(0.03), col1_w - Inches(0.3), Inches(0.3),
                'DATA SOURCES & INPUTS', font_name='Montserrat SemiBold',
                font_size=Pt(11), bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    content1_y = col_y + Inches(0.45)
    lines1 = [
        'Grid Infrastructure| screening',
        '\u2022 OSM/Overpass API: substations, lines, cables',
        '\u2022 Voltage mapping & distance scoring',
        '\u2022 User evidence ingestion (STMG, utility refs)',
        '',
        'Climate & Weather| data',
        '\u2022 Open-Meteo API (5-year archive \u2192 1 repr. year)',
        '\u2022 8,760 hourly records (T_db, RH)',
        '\u2022 Manual CSV upload supported',
        '',
        'Solar Irradiance| profile',
        '\u2022 PVGIS / Copernicus (tilt, azimuth, tech type)',
        '',
        'Site Geometry & Power|',
        '\u2022 KML/KMZ upload \u2192 area, perimeter',
        '\u2022 Power reservation + voltage levels',
        '\u2022 Buildable ratio, building height, floor count',
    ]
    add_multiline_textbox(slide, col1_x + Inches(0.12), content1_y,
                          col1_w - Inches(0.24), col_h - Inches(0.55),
                          lines1, font_size=Pt(8.5), color=DARK_TEAL,
                          bullet_color=PURPLE)

    # ── Column 2: SIMULATION ENGINE ──────────────────────────────
    box2 = add_rounded_rect(slide, col2_x, col_y, col2_w, col_h,
                            fill_color=LIGHT_GRAY, border_color=PURPLE, border_width=1.5)

    hdr2_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_x, col_y, col2_w, Inches(0.35))
    hdr2_bg.fill.solid()
    hdr2_bg.fill.fore_color.rgb = DARK_PURPLE
    no_border(hdr2_bg)
    add_textbox(slide, col2_x + Inches(0.15), col_y + Inches(0.03), col2_w - Inches(0.3), Inches(0.3),
                'SIMULATION ENGINE', font_name='Montserrat SemiBold',
                font_size=Pt(11), bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    content2_y = col_y + Inches(0.45)
    lines2 = [
        '8,760-Hour PUE Engine|',
        '\u2022 Hourly weather-driven cooling COP model',
        '\u2022 Energy-weighted annual PUE (Uptime methodology)',
        '\u2022 IT capacity spectrum: Worst / P99 / P90 / Mean / Best',
        '',
        '8 Cooling Topologies|',
        '\u2022 Air-CRAC, AHU, Chiller+Econ, RDHx, DLC, Immersion,',
        '  Dry Cooler \u2014 3-mode economizer w/ wet-bulb dispatch',
        '',
        '6 Load Types| \u00d7 4 Redundancy \u00d7 3 Density',
        '\u2022 Colocation (Std/HD), HPC, AI/GPU, Hyperscale, Edge',
        '\u2022 N / N+1 / 2N / 2N+1 redundancy levels',
        '',
        'Advanced Analytics|',
        '\u2022 Sensitivity tornado charts & break-even solver',
        '\u2022 Load-mix optimizer (multi-workload allocation)',
        '\u2022 Firm capacity solver (guaranteed MW)',
        '\u2022 PV + BESS + Fuel Cell hourly dispatch',
    ]
    add_multiline_textbox(slide, col2_x + Inches(0.12), content2_y,
                          col2_w - Inches(0.24), col_h - Inches(0.55),
                          lines2, font_size=Pt(8.5), color=DARK_TEAL,
                          bullet_color=CYAN)

    # ── Column 3: OUTPUTS & REPORTING ────────────────────────────
    box3 = add_rounded_rect(slide, col3_x, col_y, col3_w, col_h,
                            fill_color=LIGHT_GRAY, border_color=LAVENDER, border_width=1)

    hdr3_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col3_x, col_y, col3_w, Inches(0.35))
    hdr3_bg.fill.solid()
    hdr3_bg.fill.fore_color.rgb = NAVY
    no_border(hdr3_bg)
    add_textbox(slide, col3_x + Inches(0.15), col_y + Inches(0.03), col3_w - Inches(0.3), Inches(0.3),
                'OUTPUTS & REPORTING', font_name='Montserrat SemiBold',
                font_size=Pt(11), bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    content3_y = col_y + Inches(0.45)
    lines3 = [
        'Scenario Results|',
        '\u2022 IT Capacity (MW) & Facility Power',
        '\u2022 Annual PUE (static + hourly)',
        '\u2022 RAG Status (Red/Amber/Green/Blue)',
        '\u2022 Composite Score 0\u2013100',
        '\u2022 Racks deployed & binding constraint',
        '',
        'Green Energy|',
        '\u2022 PV offset & BESS cycling metrics',
        '\u2022 CO\u2082 grid offset calculation',
        '\u2022 Hourly dispatch visualization',
        '',
        'Export Formats|',
        '\u2022 PDF \u2014 executive & detailed reports',
        '\u2022 Excel \u2014 multi-sheet workbook',
        '\u2022 HTML \u2014 interactive browser report',
        '\u2022 Charts: capacity, PUE, tornado, dispatch',
    ]
    add_multiline_textbox(slide, col3_x + Inches(0.12), content3_y,
                          col3_w - Inches(0.24), col_h - Inches(0.55),
                          lines3, font_size=Pt(8.5), color=DARK_TEAL,
                          bullet_color=GREEN)

    # ── Flow arrows between columns ──────────────────────────────
    arrow_y = col_y + col_h / 2
    # Arrow from col1 to col2
    add_arrow_connector(slide,
                        col1_x + col1_w + Inches(0.02), arrow_y,
                        col2_x - Inches(0.02), arrow_y,
                        color=PURPLE, width_pt=2)
    # Arrow from col2 to col3
    add_arrow_connector(slide,
                        col2_x + col2_w + Inches(0.02), arrow_y,
                        col3_x - Inches(0.02), arrow_y,
                        color=PURPLE, width_pt=2)

    # ════════════════════════════════════════════════════════════════
    # KEY METRICS STRIP
    # ════════════════════════════════════════════════════════════════

    metrics_y = Inches(5.7)
    metrics_h = Inches(0.95)
    metric_w = Inches(2.85)
    metrics_gap = Inches(0.22)
    metrics_start_x = Inches(0.55)

    metrics = [
        ('8,760', 'Hours / Year', 'Full-year hourly simulation', PURPLE),
        ('8', 'Cooling Topologies', 'Air to immersion cooling', CYAN),
        ('6', 'Load Types', 'Colocation to AI/GPU clusters', GREEN),
        ('3', 'Export Formats', 'PDF, Excel, HTML reports', LAVENDER),
    ]

    for i, (number, label, sublabel, accent_color) in enumerate(metrics):
        x = metrics_start_x + i * (metric_w + metrics_gap)

        # Background box
        mbox = add_rounded_rect(slide, x, metrics_y, metric_w, metrics_h,
                                fill_color=WHITE, border_color=accent_color, border_width=1.5)

        # Accent bar on left
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, metrics_y, Inches(0.06), metrics_h)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        no_border(accent_bar)

        # Big number
        add_textbox(slide, x + Inches(0.2), metrics_y + Inches(0.05),
                    Inches(1.2), Inches(0.5),
                    number, font_name='Montserrat SemiBold',
                    font_size=Pt(26), bold=True, color=DARK_PURPLE,
                    alignment=PP_ALIGN.LEFT)

        # Label
        add_textbox(slide, x + Inches(0.2), metrics_y + Inches(0.48),
                    metric_w - Inches(0.35), Inches(0.22),
                    label, font_name='Open Sans',
                    font_size=Pt(10), bold=True, color=NAVY,
                    alignment=PP_ALIGN.LEFT)

        # Sublabel
        add_textbox(slide, x + Inches(0.2), metrics_y + Inches(0.67),
                    metric_w - Inches(0.35), Inches(0.2),
                    sublabel, font_name='Open Sans',
                    font_size=Pt(7.5), bold=False, color=BODY_GRAY,
                    alignment=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════════════
    # FOOTER
    # ════════════════════════════════════════════════════════════════

    # Bottom bar
    bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), SLIDE_W, Inches(0.35))
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = NAVY
    no_border(bot_bar)

    add_textbox(slide, Inches(0.5), Inches(7.17), Inches(5), Inches(0.3),
                'Metlen  \u00b7  Data Center BU  \u00b7  Innovative Team  \u00b7  2026',
                font_name='Open Sans', font_size=Pt(9), bold=False, color=WHITE,
                alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(9), Inches(7.17), Inches(4), Inches(0.3),
                'DC Feasibility Tool v4.1',
                font_name='Open Sans', font_size=Pt(9), bold=False, color=LAVENDER,
                alignment=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════════════════════════
    # SAVE
    # ════════════════════════════════════════════════════════════════

    output_path = os.path.join(os.path.dirname(__file__), 'dc_feasibility_v4_intro.pptx')
    prs.save(output_path)
    print(f'Saved: {output_path}')


if __name__ == '__main__':
    build_slide()
